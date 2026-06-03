from __future__ import annotations

import csv
import re
import tempfile
from pathlib import Path
from typing import Protocol

from .models import GenerationConfig, MaterialDocument


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".pptx",
    ".docx",
    ".txt",
    ".md",
    ".xlsx",
    ".csv",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
GENERATED_OUTPUT_DIR = re.compile(r"^\d{8}_\d{6}_.+")
REFERENCE_PATH_PARTS = {
    "reference",
    "references",
    "textbook",
    "textbooks",
    "book",
    "books",
    "reading",
    "readings",
    "supplement",
    "supplementary",
    "参考",
    "教材",
    "教科书",
}
REFERENCE_PDF_PAGE_THRESHOLD = 120
REFERENCE_PDF_EXTRACT_PAGES = 30


class VisionClient(Protocol):
    def ocr_image(self, image_path: Path) -> str:
        ...


def scan_materials(input_dir: Path) -> list[MaterialDocument]:
    root = Path(input_dir).expanduser().resolve()
    if not root.exists() or not root.is_dir():
        raise FileNotFoundError(f"Input folder does not exist: {root}")

    documents: list[MaterialDocument] = []
    for file_path in sorted(p for p in root.rglob("*") if p.is_file()):
        relative_parts = file_path.relative_to(root).parts
        if any(GENERATED_OUTPUT_DIR.match(part) for part in relative_parts):
            continue
        ext = file_path.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        documents.append(
            MaterialDocument(
                source_path=file_path,
                relative_path=file_path.relative_to(root).as_posix(),
                material_type=ext.lstrip("."),
                role=_initial_material_role(file_path.relative_to(root)),
                status="indexed",
            )
        )
    return documents


def extract_materials(
    documents: list[MaterialDocument],
    client: VisionClient | None,
    config: GenerationConfig,
) -> list[MaterialDocument]:
    extracted: list[MaterialDocument] = []
    for doc in documents:
        try:
            ext = doc.source_path.suffix.lower()
            if ext in {".txt", ".md"}:
                doc.text = _read_text(doc.source_path)
            elif ext == ".csv":
                doc.text = _read_csv(doc.source_path)
            elif ext == ".docx":
                doc.text = _read_docx(doc.source_path)
            elif ext == ".pptx":
                doc.text = _read_pptx(doc.source_path)
            elif ext == ".xlsx":
                doc.text = _read_xlsx(doc.source_path)
            elif ext == ".pdf":
                if doc.role == "reference" or _is_large_reference_pdf(doc.source_path):
                    doc.role = "reference"
                    doc.text = _read_pdf(doc.source_path, max_pages=REFERENCE_PDF_EXTRACT_PAGES)
                    doc.warnings.append(
                        f"Reference PDF: extracted first {REFERENCE_PDF_EXTRACT_PAGES} pages only; "
                        "use as supporting context, not primary coverage material."
                    )
                else:
                    doc.text = _read_pdf(doc.source_path)
                if not doc.text.strip() and config.enable_vision_ocr:
                    doc.text = _ocr_pdf_pages(doc.source_path, client)
            elif ext in IMAGE_EXTENSIONS:
                doc.image_path = doc.source_path
                doc.text = _ocr_image(doc.source_path, client, config.enable_vision_ocr)
            doc.status = "extracted" if doc.text.strip() else "empty"
            if not doc.text.strip():
                doc.warnings.append("No text extracted; check whether the file is scanned or unsupported.")
        except Exception as exc:  # noqa: BLE001 - keep batch processing alive.
            doc.status = "failed"
            doc.warnings.append(str(exc))
        extracted.append(doc)
    return extracted


def _initial_material_role(relative_path: Path) -> str:
    parts = {part.casefold() for part in relative_path.parts}
    stem = relative_path.stem.casefold()
    if parts & REFERENCE_PATH_PARTS:
        return "reference"
    if any(marker in stem for marker in REFERENCE_PATH_PARTS):
        return "reference"
    return "primary"


def _is_large_reference_pdf(path: Path) -> bool:
    if path.suffix.lower() != ".pdf":
        return False
    try:
        from pypdf import PdfReader
    except ImportError:
        return False
    try:
        reader = PdfReader(str(path))
        return len(reader.pages) >= REFERENCE_PDF_PAGE_THRESHOLD
    except Exception:
        return False


def _read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(errors="replace")


def _read_csv(path: Path) -> str:
    rows: list[str] = []
    with path.open("r", encoding="utf-8-sig", newline="", errors="replace") as fh:
        reader = csv.reader(fh)
        for idx, row in enumerate(reader, start=1):
            rows.append(f"Row {idx}: " + " | ".join(row))
    return "\n".join(rows)


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to read DOCX files.") from exc
    document = Document(path)
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        paragraphs.append(f"[Table {table_index}]")
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(paragraphs)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to read PPTX files.") from exc
    prs = Presentation(path)
    slides: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = [f"[Slide {index}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _read_xlsx(path: Path) -> str:
    try:
        import pandas as pd
    except ImportError as exc:
        raise RuntimeError("pandas and openpyxl are required to read XLSX files.") from exc
    workbook = pd.read_excel(path, sheet_name=None, dtype=str)
    parts: list[str] = []
    for sheet_name, frame in workbook.items():
        parts.append(f"[Sheet: {sheet_name}]")
        parts.append(frame.fillna("").to_csv(index=False))
    return "\n".join(parts)


def _read_pdf(path: Path, max_pages: int | None = None) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is required to read PDF files.") from exc
    reader = PdfReader(str(path))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        if max_pages is not None and index > max_pages:
            break
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {index}]\n{text}")
    return "\n\n".join(pages)


def _ocr_image(path: Path, client: VisionClient | None, enabled: bool) -> str:
    if not enabled:
        return ""
    if client is None:
        raise RuntimeError("Vision OCR needs a configured API client.")
    return client.ocr_image(path)


def _ocr_pdf_pages(path: Path, client: VisionClient | None) -> str:
    if client is None:
        raise RuntimeError("Vision OCR needs a configured API client.")
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required to OCR scanned PDF pages.") from exc

    chunks: list[str] = []
    with tempfile.TemporaryDirectory() as temp_dir:
        pdf = fitz.open(path)
        for page_index in range(len(pdf)):
            page = pdf.load_page(page_index)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = Path(temp_dir) / f"page_{page_index + 1}.png"
            pix.save(image_path)
            text = client.ocr_image(image_path)
            chunks.append(f"[Page {page_index + 1} OCR]\n{text}")
    return "\n\n".join(chunks)
